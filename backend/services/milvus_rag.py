import os
import logging
from typing import List, Dict, Any, Optional
from pathlib import Path
import json
import numpy as np
from pymilvus import MilvusClient, DataType, Collection
import openai
from openai import OpenAI
import hashlib

logger = logging.getLogger(__name__)

class MilvusRAG:
    def __init__(self):
        self.client = None
        self.openai_client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"), timeout=float(os.getenv("OPENAI_TIMEOUT", "180")))
        self.db_path = "./milvus_rfp.db"
        self.collection_name = "kamiwaza_knowledge"
        self.embedding_model = "text-embedding-3-large"
        self.dimension = 3072  # Dimension for text-embedding-3-large
        
    async def initialize(self):
        """Initialize Milvus Lite database and create collections"""
        try:
            self.client = MilvusClient(uri=self.db_path)
            
            if not self.client.has_collection(self.collection_name):
                self.client.create_collection(
                    collection_name=self.collection_name,
                    dimension=self.dimension,
                    metric_type="COSINE",
                    consistency_level="Strong"
                )
                logger.info(f"Created collection: {self.collection_name}")
            
            await self._initialize_company_knowledge()
            logger.info("Milvus RAG initialized successfully")
            
        except Exception as e:
            logger.error(f"Failed to initialize Milvus: {e}")
            raise
    
    async def _initialize_company_knowledge(self):
        """Load initial Kamiwaza company knowledge"""
        kamiwaza_info = [
            {
                "content": """Kamiwaza is a cutting-edge technology company specializing in AI-powered solutions, 
                cloud infrastructure, and digital transformation services. We have extensive experience in 
                government contracting, particularly with federal agencies.""",
                "metadata": {"type": "company_overview", "category": "general"}
            },
            {
                "content": """Core Competencies:
                - Artificial Intelligence and Machine Learning Solutions
                - Cloud Migration and Infrastructure (AWS, Azure, GCP certified)
                - Cybersecurity and Risk Management
                - Data Analytics and Business Intelligence
                - Custom Software Development
                - DevOps and Automation
                - Digital Transformation Consulting""",
                "metadata": {"type": "capabilities", "category": "technical"}
            },
            {
                "content": """Past Performance:
                - Successfully delivered 50+ federal contracts
                - 98% customer satisfaction rating
                - On-time delivery rate: 100%
                - Budget compliance: 100%
                - Security clearance: Up to Secret level
                - CMMI Level 3 certified
                - ISO 9001:2015 certified""",
                "metadata": {"type": "past_performance", "category": "qualifications"}
            },
            {
                "content": """Certifications and Compliance:
                - Small Business (SB) certified
                - FedRAMP authorized
                - NIST 800-171 compliant
                - SOC 2 Type II certified
                - HIPAA compliant
                - Section 508 accessibility compliant""",
                "metadata": {"type": "certifications", "category": "compliance"}
            },
            {
                "content": """Key Personnel:
                - CEO: 20+ years federal contracting experience
                - CTO: Former DoD technology advisor
                - VP Engineering: Led teams at Fortune 500 tech companies
                - VP Sales: 15+ years government sales experience
                - Security Officer: Former NSA cybersecurity expert""",
                "metadata": {"type": "personnel", "category": "team"}
            },
            {
                "content": """Contract Vehicles:
                - GSA Schedule 70
                - CIO-SP3 Small Business
                - SEWP V
                - 8(a) STARS III
                - Alliant 2 Small Business""",
                "metadata": {"type": "contract_vehicles", "category": "contracting"}
            },
            {
                "content": """Technical Stack and Tools:
                - Languages: Python, Java, JavaScript, Go, C++
                - Frameworks: React, Angular, Django, Spring Boot, .NET
                - Databases: PostgreSQL, MongoDB, Oracle, Elasticsearch
                - Cloud: AWS (Advanced Partner), Azure (Gold Partner), GCP
                - AI/ML: TensorFlow, PyTorch, OpenAI, Hugging Face
                - DevOps: Kubernetes, Docker, Jenkins, GitLab CI/CD""",
                "metadata": {"type": "tech_stack", "category": "technical"}
            },
            {
                "content": """Differentiators:
                - Rapid deployment capabilities (2-4 week implementation)
                - 24/7 US-based support team
                - Proprietary AI acceleration framework
                - Zero security incidents in company history
                - Agile development methodology with 2-week sprints
                - Cost-effective solutions averaging 20% below industry rates""",
                "metadata": {"type": "differentiators", "category": "competitive_advantage"}
            }
        ]
        
        await self.index_documents(kamiwaza_info)
    
    async def index_documents(self, documents: List[Dict[str, Any]]):
        """Index documents into Milvus"""
        try:
            data_to_insert = []
            
            for idx, doc in enumerate(documents):
                content = doc.get("content", "")
                metadata = doc.get("metadata", {})
                
                # Use hash to create a unique integer ID
                doc_hash = hashlib.md5(content.encode()).hexdigest()
                # Convert first 16 chars of hash to int
                doc_id = int(doc_hash[:16], 16) % (2**63 - 1)  # Ensure it fits in int64
                
                embedding = await self.get_embedding(content)
                
                data_to_insert.append({
                    "id": doc_id,
                    "vector": embedding,
                    "text": content,
                    "metadata": json.dumps(metadata)
                })
            
            if data_to_insert:
                self.client.insert(
                    collection_name=self.collection_name,
                    data=data_to_insert
                )
                logger.info(f"Indexed {len(data_to_insert)} documents")
                
        except Exception as e:
            logger.error(f"Error indexing documents: {e}")
            raise
    
    async def get_embedding(self, text: str) -> List[float]:
        """Get embedding for text using OpenAI"""
        try:
            response = self.openai_client.embeddings.create(
                input=text,
                model=self.embedding_model
            )
            return response.data[0].embedding
        except Exception as e:
            logger.error(f"Error getting embedding: {e}")
            raise
    
    async def search_similar(self, query: str, limit: int = 5) -> List[Dict[str, Any]]:
        """Search for similar documents"""
        try:
            query_embedding = await self.get_embedding(query)
            
            results = self.client.search(
                collection_name=self.collection_name,
                data=[query_embedding],
                limit=limit,
                output_fields=["text", "metadata"]
            )
            
            formatted_results = []
            for hits in results:
                for hit in hits:
                    formatted_results.append({
                        "content": hit.get("entity", {}).get("text", ""),
                        "metadata": json.loads(hit.get("entity", {}).get("metadata", "{}")),
                        "distance": hit.get("distance", 1.0)
                    })
            
            return formatted_results
            
        except Exception as e:
            logger.error(f"Error searching documents: {e}")
            return []
    
    async def get_company_context(self, query: str = "Kamiwaza company information") -> str:
        """Get relevant company context for RFP response"""
        try:
            results = await self.search_similar(query, limit=10)
            
            context = "KAMIWAZA COMPANY INFORMATION:\n\n"
            for result in results:
                context += f"{result['content']}\n\n"
            
            return context
            
        except Exception as e:
            logger.error(f"Error getting company context: {e}")
            return "Kamiwaza is a technology company specializing in AI and cloud solutions."
    
    async def add_rfp_response(self, rfp_id: str, response_data: Dict[str, Any]):
        """Store completed RFP response for future reference"""
        try:
            documents = []
            
            for section_name, content in response_data.items():
                if isinstance(content, str):
                    documents.append({
                        "content": content,
                        "metadata": {
                            "type": "rfp_response",
                            "rfp_id": rfp_id,
                            "section": section_name,
                            "date": str(datetime.now())
                        }
                    })
            
            await self.index_documents(documents)
            logger.info(f"Stored RFP response for {rfp_id}")
            
        except Exception as e:
            logger.error(f"Error storing RFP response: {e}")
    
    async def search_past_rfps(self, requirements: str) -> List[Dict[str, Any]]:
        """Search for similar past RFP responses"""
        try:
            results = await self.search_similar(requirements, limit=5)
            
            past_rfps = [r for r in results if r.get("metadata", {}).get("type") == "rfp_response"]
            
            return past_rfps
            
        except Exception as e:
            logger.error(f"Error searching past RFPs: {e}")
            return []
    
    async def index_company_documents(self, folder_path: Optional[str] = None) -> Dict[str, Any]:
        """Index company documents from a folder"""
        try:
            if not folder_path:
                folder_path = "./company_documents"
            
            folder = Path(folder_path)
            if not folder.exists():
                logger.warning(f"Folder {folder_path} does not exist")
                return {"indexed": 0, "message": "Folder not found"}
            
            documents = []
            for file_path in folder.glob("**/*"):
                if file_path.is_file() and file_path.suffix in ['.txt', '.md', '.pdf', '.docx']:
                    try:
                        content = self._read_file(file_path)
                        if content:
                            documents.append({
                                "content": content,
                                "metadata": {
                                    "type": "company_document",
                                    "filename": file_path.name,
                                    "path": str(file_path)
                                }
                            })
                    except Exception as e:
                        logger.warning(f"Could not read file {file_path}: {e}")
            
            if documents:
                await self.index_documents(documents)
            
            return {
                "indexed": len(documents),
                "message": f"Successfully indexed {len(documents)} documents"
            }
            
        except Exception as e:
            logger.error(f"Error indexing company documents: {e}")
            raise
    
    def _read_file(self, file_path: Path) -> Optional[str]:
        """Read content from various file types"""
        try:
            suffix_lower = file_path.suffix.lower()
            if suffix_lower in ['.txt', '.md']:
                return file_path.read_text(encoding='utf-8')
            elif suffix_lower == '.pdf':
                import pypdf
                reader = pypdf.PdfReader(str(file_path))
                text = ""
                for page in reader.pages:
                    text += page.extract_text()
                return text
            elif suffix_lower == '.docx':
                import docx
                doc = docx.Document(str(file_path))
                return "\n".join([para.text for para in doc.paragraphs])
            elif suffix_lower in ['.xlsx', '.xls']:
                import pandas as pd
                df = pd.read_excel(str(file_path), sheet_name=None)
                text = ""
                for sheet_name, sheet_df in df.items():
                    text += f"\n=== Sheet: {sheet_name} ===\n"
                    text += sheet_df.to_string()
                return text
            elif suffix_lower == '.rtf':
                # For RTF files, use striprtf library or basic extraction
                try:
                    from striprtf.striprtf import rtf_to_text
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        rtf_content = f.read()
                    return rtf_to_text(rtf_content)
                except ImportError:
                    # Fallback: basic RTF text extraction
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        content = f.read()
                    # Remove basic RTF formatting
                    import re
                    # Remove RTF control words
                    content = re.sub(r'\\[a-z]+\d*\s?', '', content)
                    # Remove braces
                    content = re.sub(r'[{}]', '', content)
                    return content
            elif suffix_lower == '.pptx':
                from pptx import Presentation
                prs = Presentation(str(file_path))
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text
            else:
                return None
        except Exception as e:
            logger.error(f"Error reading file {file_path}: {e}")
            return None

    async def add_refinement_session(self, session_data: Dict[str, Any]) -> bool:
        """Save refinement session data to RAG"""
        try:
            # Convert session data to documents for indexing
            documents_to_index = []

            # Index each Q&A pair
            for qa in session_data.get("questions_answered", []):
                doc_content = f"""
                Refinement Knowledge Entry:
                Question: {qa.get('question', '')}
                Human Answer: {qa.get('answer', '')}
                Context: {session_data.get('notice_id', 'Unknown')}
                Timestamp: {session_data.get('timestamp', '')}
                """
                documents_to_index.append(doc_content)

            # Index refined documents snippets
            for doc_name, content in session_data.get("refined_documents", {}).items():
                # Just index key improvements, not full documents
                if len(content) > 1000:
                    content = content[:1000] + "..."
                documents_to_index.append(f"Refined {doc_name}: {content}")

            # Index documents if we have any
            if documents_to_index:
                await self.index_documents(documents_to_index)
                logger.info(f"Indexed {len(documents_to_index)} refinement entries")

            return True
        except Exception as e:
            logger.error(f"Error adding refinement session to RAG: {e}")
            return False

from datetime import datetime
